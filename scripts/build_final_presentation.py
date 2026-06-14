from pathlib import Path

import matplotlib.pyplot as plt
import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DATA_PATH = ROOT / "data" / "raw" / "sample_ocean_data.csv"
SCREENSHOT_PATH = ROOT / "docs" / "screenshots" / "ocean_dashboard_home.png"
ASSET_DIR = ROOT / "docs" / "presentation_assets"
OUTPUT_PATH = ROOT / "Ocean_Environment_Dashboard_Final_Presentation.pptx"

FONT = "Microsoft JhengHei"
FONT_EN = "Arial"

COLORS = {
    "navy": "073B4C",
    "blue": "0077B6",
    "cyan": "00B4D8",
    "aqua": "90E0EF",
    "pale": "EAF8FB",
    "white": "FFFFFF",
    "gray": "EEF2F5",
    "text": "12313F",
    "muted": "5A6B75",
    "green": "2A9D8F",
    "orange": "F4A261",
}

STATION_ZH = {
    "Keelung": "基隆",
    "Taichung": "台中",
    "Kaohsiung": "高雄",
    "Hualien": "花蓮",
    "Taitung": "台東",
}


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def set_background(slide, color="FFFFFF"):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(color)


def add_textbox(slide, text, x, y, w, h, font_size=18, color="text", bold=False,
                align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, font_face=FONT,
                line_spacing=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.03)
    frame.margin_right = Inches(0.03)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    if line_spacing:
        paragraph.line_spacing = line_spacing
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_face
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = rgb(COLORS.get(color, color))
    return box


def add_bullets(slide, items, x, y, w, h, font_size=17, color="text",
                bullet_color="blue", line_gap=0.38):
    for index, item in enumerate(items):
        yy = y + index * line_gap
        slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x),
            Inches(yy + 0.07),
            Inches(0.12),
            Inches(0.12),
        ).fill.solid()
        dot = slide.shapes[-1]
        dot.fill.fore_color.rgb = rgb(COLORS[bullet_color])
        dot.line.fill.background()
        add_textbox(slide, item, x + 0.22, yy, w - 0.22, 0.32, font_size, color)


def add_page_title(slide, title, subtitle=None, page=None):
    add_textbox(slide, title, 0.48, 0.32, 6.6, 0.46, 24, "navy", True)
    slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.48), Inches(0.86), Inches(1.05), Inches(0.05)
    ).fill.solid()
    line = slide.shapes[-1]
    line.fill.fore_color.rgb = rgb(COLORS["cyan"])
    line.line.fill.background()
    if subtitle:
        add_textbox(slide, subtitle, 0.48, 0.98, 7.8, 0.32, 11.5, "muted")
    if page is not None:
        add_page_badge(slide, page)


def add_page_badge(slide, page):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.28), Inches(5.08), Inches(0.43), Inches(0.34)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(COLORS["blue"])
    shape.line.fill.background()
    add_textbox(slide, str(page), 9.28, 5.105, 0.43, 0.18, 9, "white", True, PP_ALIGN.CENTER)


def add_card(slide, x, y, w, h, title, body, accent="cyan", title_size=15, body_size=12):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = rgb(COLORS["white"])
    card.line.color.rgb = rgb(COLORS["aqua"])
    card.line.width = Pt(1)
    marker = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    marker.fill.solid()
    marker.fill.fore_color.rgb = rgb(COLORS[accent])
    marker.line.fill.background()
    add_textbox(slide, title, x + 0.18, y + 0.14, w - 0.28, 0.3, title_size, "navy", True)
    add_textbox(slide, body, x + 0.18, y + 0.52, w - 0.28, h - 0.62, body_size, "text")


def add_badge(slide, text, x, y, w=0.72, color="blue"):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.28))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(COLORS[color])
    shape.line.fill.background()
    add_textbox(slide, text, x, y + 0.045, w, 0.14, 8.5, "white", True, PP_ALIGN.CENTER)


def add_picture_fit(slide, image_path, x, y, w, h):
    image = Image.open(image_path)
    img_w, img_h = image.size
    box_ratio = w / h
    img_ratio = img_w / img_h
    if img_ratio > box_ratio:
        final_w = w
        final_h = w / img_ratio
        final_x = x
        final_y = y + (h - final_h) / 2
    else:
        final_h = h
        final_w = h * img_ratio
        final_x = x + (w - final_w) / 2
        final_y = y
    return slide.shapes.add_picture(str(image_path), Inches(final_x), Inches(final_y), Inches(final_w), Inches(final_h))


def make_chart_assets():
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    df = pd.read_csv(DATA_PATH)
    df["date"] = pd.to_datetime(df["date"])
    df["station_zh"] = df["station"].map(STATION_ZH)

    plt.rcParams["font.sans-serif"] = ["Microsoft JhengHei", "Noto Sans CJK TC", "SimHei", "Arial Unicode MS"]
    plt.rcParams["axes.unicode_minus"] = False

    chart_paths = {}

    fig, ax = plt.subplots(figsize=(9, 4.8), dpi=160)
    for station, station_df in df.groupby("station_zh"):
        ax.plot(station_df["date"], station_df["sea_temperature_c"], marker="o", linewidth=2, label=station)
    ax.set_title("Sea Temperature Trend", fontsize=15, weight="bold")
    ax.set_xlabel("Date")
    ax.set_ylabel("Sea Temperature (°C)")
    ax.grid(alpha=0.25)
    ax.legend(fontsize=9)
    fig.autofmt_xdate(rotation=25)
    fig.tight_layout()
    chart_paths["temperature"] = ASSET_DIR / "sea_temperature_trend.png"
    fig.savefig(chart_paths["temperature"], transparent=False, facecolor="white")
    plt.close(fig)

    fig, ax = plt.subplots(figsize=(8.2, 4.6), dpi=160)
    station_avg = df.groupby("station_zh", as_index=False)["sea_temperature_c"].mean().sort_values("sea_temperature_c")
    ax.bar(station_avg["station_zh"], station_avg["sea_temperature_c"], color="#00A6A6", edgecolor="#0077B6")
    ax.set_title("Average Sea Temperature by Station", fontsize=15, weight="bold")
    ax.set_xlabel("Station")
    ax.set_ylabel("Average Sea Temperature (°C)")
    ax.grid(axis="y", alpha=0.25)
    fig.tight_layout()
    chart_paths["bar"] = ASSET_DIR / "station_average_temperature.png"
    fig.savefig(chart_paths["bar"], transparent=False, facecolor="white")
    plt.close(fig)

    fig, ax = plt.subplots(figsize=(8.2, 4.6), dpi=160)
    for station, station_df in df.groupby("station_zh"):
        ax.plot(station_df["date"], station_df["wave_height_m"], marker="o", linewidth=2, label=station)
    ax.set_title("Wave Height Trend", fontsize=15, weight="bold")
    ax.set_xlabel("Date")
    ax.set_ylabel("Wave Height (m)")
    ax.grid(alpha=0.25)
    ax.legend(fontsize=8)
    fig.autofmt_xdate(rotation=25)
    fig.tight_layout()
    chart_paths["wave"] = ASSET_DIR / "wave_height_trend.png"
    fig.savefig(chart_paths["wave"], transparent=False, facecolor="white")
    plt.close(fig)

    sample = df.head(8)[["date", "station", "sea_temperature_c", "tide_level_m", "wave_height_m", "wind_speed_mps", "salinity_psu"]].copy()
    sample["date"] = sample["date"].dt.strftime("%Y-%m-%d")
    sample["station"] = sample["station"].map(STATION_ZH)
    fig, ax = plt.subplots(figsize=(9, 2.6), dpi=160)
    ax.axis("off")
    table = ax.table(
        cellText=sample.values,
        colLabels=["日期", "測站", "海溫", "潮位", "浪高", "風速", "鹽度"],
        cellLoc="center",
        loc="center",
    )
    table.auto_set_font_size(False)
    table.set_fontsize(8.5)
    table.scale(1, 1.35)
    for (row, col), cell in table.get_celld().items():
        cell.set_edgecolor("#CDEAF2")
        if row == 0:
            cell.set_facecolor("#0077B6")
            cell.get_text().set_color("white")
            cell.get_text().set_weight("bold")
        else:
            cell.set_facecolor("#F6FBFF" if row % 2 else "white")
    fig.tight_layout()
    chart_paths["table"] = ASSET_DIR / "sample_data_table.png"
    fig.savefig(chart_paths["table"], transparent=False, facecolor="white")
    plt.close(fig)

    return chart_paths


def add_cover(slide):
    set_background(slide, COLORS["navy"])
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(5.625)).fill.solid()
    bg = slide.shapes[-1]
    bg.fill.fore_color.rgb = rgb(COLORS["navy"])
    bg.line.fill.background()
    for idx, color in enumerate(["0077B6", "00B4D8", "90E0EF"]):
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-0.5 + idx * 1.8), Inches(4.55 - idx * 0.12), Inches(4.2), Inches(1.15))
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(color)
        shape.line.fill.background()
    add_textbox(slide, "海洋環境資料儀表板", 0.75, 1.55, 8.5, 0.75, 34, "white", True, PP_ALIGN.CENTER)
    add_textbox(slide, "114-2 巨量資料與雲端運算期末專題", 1.45, 2.35, 7.1, 0.38, 18, "aqua", False, PP_ALIGN.CENTER)
    add_textbox(slide, "小組專題", 4.3, 3.0, 1.4, 0.28, 13, "white", False, PP_ALIGN.CENTER)
    for i, tech in enumerate(["Python", "Pandas", "Matplotlib", "Gradio", "Docker"]):
        add_badge(slide, tech, 2.25 + i * 1.15, 3.55, 0.9, "blue" if i % 2 == 0 else "green")
    add_textbox(
        slide,
        "組員 1：________  學號：________    組員 2：________  學號：________\n"
        "組員 3：________  學號：________    組員 4：________  學號：________\n"
        "組員 5：________  學號：________",
        2.15,
        4.42,
        5.7,
        0.62,
        8.5,
        "white",
        False,
        PP_ALIGN.CENTER,
    )


def create_presentation(chart_paths):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    blank = prs.slide_layouts[6]

    # 1
    slide = prs.slides.add_slide(blank)
    add_cover(slide)

    # 2
    slide = prs.slides.add_slide(blank)
    set_background(slide, COLORS["pale"])
    add_page_title(slide, "報告大綱", "以四個部分說明專題從背景到成果的完整流程", 2)
    agenda = [
        ("01", "背景介紹", "專題背景、動機與目的"),
        ("02", "系統流程與方法", "架構、資料處理與使用技術"),
        ("03", "結果與展示", "資料視覺化、Gradio 與 Docker"),
        ("04", "結論", "專題成果與未來改善方向"),
    ]
    for i, (num, title, body) in enumerate(agenda):
        x = 0.7 + (i % 2) * 4.55
        y = 1.35 + (i // 2) * 1.45
        add_card(slide, x, y, 4.05, 1.05, f"{num}  {title}", body, ["blue", "green", "cyan", "orange"][i], 17, 13)

    # 3
    slide = prs.slides.add_slide(blank)
    set_background(slide, COLORS["white"])
    add_page_title(slide, "背景介紹", "海洋環境資料能協助觀察不同海域與時間的變化", 3)
    add_bullets(slide, [
        "海洋環境資料包含海溫、潮位、浪高、風速、鹽度等資訊。",
        "這些資料可用於觀察不同海域的環境變化。",
        "對海事活動、港口管理、漁業活動與海洋觀測都有參考價值。",
        "透過儀表板呈現資料，能讓使用者更快速理解資料趨勢。",
    ], 0.72, 1.32, 5.1, 2.0, 15.5)
    for i, label in enumerate(["海象觀測", "港口管理", "資料趨勢"]):
        add_card(slide, 6.25, 1.15 + i * 1.05, 2.8, 0.78, label, ["蒐集多測站環境資訊", "支援海域活動判斷", "用圖表快速理解變化"][i], ["blue", "green", "orange"][i], 14, 10.5)

    # 4
    slide = prs.slides.add_slide(blank)
    set_background(slide, COLORS["pale"])
    add_page_title(slide, "專題動機", "將表格資料轉換成可互動分析的資料應用系統", 4)
    add_picture_fit(slide, chart_paths["table"], 0.55, 1.28, 3.8, 1.35)
    add_textbox(slide, "原始 CSV / 表格", 1.28, 2.7, 2.3, 0.28, 13, "navy", True, PP_ALIGN.CENTER)
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.42), Inches(1.78), Inches(1.12), Inches(0.42))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = rgb(COLORS["cyan"])
    arrow.line.fill.background()
    add_picture_fit(slide, SCREENSHOT_PATH, 5.78, 1.05, 3.55, 2.02)
    add_textbox(slide, "互動式儀表板", 6.35, 3.13, 2.4, 0.28, 13, "navy", True, PP_ALIGN.CENTER)
    add_textbox(slide, "重點句", 0.75, 3.75, 0.9, 0.25, 12, "white", True, PP_ALIGN.CENTER)
    tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(3.7), Inches(0.98), Inches(0.32))
    tag.fill.solid()
    tag.fill.fore_color.rgb = rgb(COLORS["blue"])
    tag.line.fill.background()
    add_textbox(slide, "將原始海洋資料轉換成可互動分析、可視覺化呈現、可 Docker 部署的資料應用系統。", 1.8, 3.72, 7.3, 0.42, 16, "text", True)

    # 5
    slide = prs.slides.add_slide(blank)
    set_background(slide, COLORS["white"])
    add_page_title(slide, "專題目的", "以三個目標串起資料處理、分析視覺化與部署", 5)
    purposes = [
        ("建立資料清洗流程", "使用 Pandas 處理日期格式、缺失值、重複資料與數值欄位。", "blue"),
        ("建立分析與視覺化功能", "使用 NumPy、Pandas、Matplotlib 產生統計摘要與圖表。", "green"),
        ("建立可部署互動系統", "使用 Gradio 建立網頁介面，並透過 Docker 容器化部署。", "orange"),
    ]
    for i, (title, body, color) in enumerate(purposes):
        add_card(slide, 0.65 + i * 3.05, 1.55, 2.65, 2.15, title, body, color, 15, 12.2)
        add_badge(slide, f"0{i+1}", 1.55 + i * 3.05, 3.92, 0.6, color)

    # 6
    slide = prs.slides.add_slide(blank)
    set_background(slide, COLORS["pale"])
    add_page_title(slide, "系統架構", "CSV 資料經過清洗、分析與視覺化後，以 Gradio + Docker 展示", 6)
    stages = ["CSV 資料", "資料讀取", "資料清洗", "統計分析", "視覺化圖表", "Gradio 儀表板", "Docker 部署"]
    for i, stage in enumerate(stages):
        x = 0.38 + i * 1.34
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.75), Inches(1.05), Inches(0.62))
        box.fill.solid()
        box.fill.fore_color.rgb = rgb(COLORS["white"])
        box.line.color.rgb = rgb(COLORS["cyan"])
        add_textbox(slide, stage, x + 0.05, 1.93, 0.95, 0.16, 9.5, "navy", True, PP_ALIGN.CENTER)
        if i < len(stages) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 1.08), Inches(1.93), Inches(0.22), Inches(0.18))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = rgb(COLORS["blue"])
            arrow.line.fill.background()
    techs = ["Python", "Pandas", "NumPy", "Matplotlib", "Gradio", "Docker", "GitHub"]
    for i, tech in enumerate(techs):
        add_badge(slide, tech, 1.05 + (i % 4) * 1.8, 3.25 + (i // 4) * 0.48, 1.05, ["blue", "green", "cyan", "orange"][i % 4])

    # 7
    slide = prs.slides.add_slide(blank)
    set_background(slide, COLORS["white"])
    add_page_title(slide, "資料來源與欄位說明", "使用 sample_ocean_data.csv 作為展示資料", 7)
    add_picture_fit(slide, chart_paths["table"], 0.55, 1.12, 4.25, 1.65)
    fields = [
        ("date", "日期"), ("station", "測站名稱"), ("sea_temperature_c", "海溫"),
        ("tide_level_m", "潮位"), ("wave_height_m", "浪高"),
        ("wind_speed_mps", "風速"), ("salinity_psu", "鹽度"),
    ]
    y = 1.1
    for idx, (field, desc) in enumerate(fields):
        add_textbox(slide, field, 5.15, y + idx * 0.42, 2.2, 0.22, 11, "navy", True, font_face=FONT_EN)
        add_textbox(slide, desc, 7.45, y + idx * 0.42, 1.6, 0.22, 11, "text")
    add_textbox(slide, "測站：Keelung、Taichung、Kaohsiung、Hualien、Taitung", 0.8, 4.42, 8.3, 0.3, 14, "muted", False, PP_ALIGN.CENTER)

    # 8
    slide = prs.slides.add_slide(blank)
    set_background(slide, COLORS["pale"])
    add_page_title(slide, "資料清洗流程", "分析前先讓資料格式一致，避免統計結果失真", 8)
    flow = [("原始資料", "CSV"), ("清洗規則", "Pandas"), ("乾淨資料", "Processed")]
    for i, (title, sub) in enumerate(flow):
        add_card(slide, 0.9 + i * 3.05, 1.28, 2.25, 0.88, title, sub, ["blue", "green", "cyan"][i], 16, 12)
        if i < 2:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.25 + i * 3.05), Inches(1.58), Inches(0.45), Inches(0.22))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = rgb(COLORS["blue"])
            arrow.line.fill.background()
    add_bullets(slide, [
        "日期格式轉換：將 date 欄位轉成 datetime。",
        "數值欄位轉換：海溫、潮位、浪高、風速、鹽度轉成 numeric。",
        "缺失值處理：數值欄位用平均值補值，測站缺失標記 Unknown。",
        "重複資料處理：移除重複列，避免統計失真。",
        "資料排序：依日期排序，方便趨勢分析。",
    ], 0.95, 2.75, 8.1, 2.1, 13.5, "text", "cyan", 0.36)

    # 9
    slide = prs.slides.add_slide(blank)
    set_background(slide, COLORS["white"])
    add_page_title(slide, "資料分析方法", "以統計摘要、測站分組、月份平均與趨勢圖掌握資料", 9)
    methods = [
        ("整體統計摘要", "平均值、最大值、最小值、標準差"),
        ("測站分組統計", "比較不同測站的各項環境差異"),
        ("月份平均分析", "觀察不同月份的海洋環境變化"),
        ("趨勢視覺化", "使用折線圖與長條圖呈現資料變化"),
    ]
    for i, (title, body) in enumerate(methods):
        add_card(slide, 0.7 + (i % 2) * 4.25, 1.25 + (i // 2) * 1.45, 3.8, 1.05, title, body, ["blue", "green", "cyan", "orange"][i], 15, 12)

    # 10
    slide = prs.slides.add_slide(blank)
    set_background(slide, COLORS["pale"])
    add_page_title(slide, "視覺化結果", "透過圖表快速觀察不同測站與不同時間的資料變化", 10)
    add_picture_fit(slide, chart_paths["temperature"], 0.45, 1.1, 4.35, 2.1)
    add_picture_fit(slide, chart_paths["bar"], 5.05, 1.1, 4.35, 2.1)
    add_picture_fit(slide, chart_paths["wave"], 2.75, 3.25, 4.5, 1.72)

    # 11
    slide = prs.slides.add_slide(blank)
    set_background(slide, COLORS["white"])
    add_page_title(slide, "系統展示與 Docker 部署", "Gradio 介面負責互動操作，Docker 負責環境重現", 11)
    add_picture_fit(slide, SCREENSHOT_PATH, 0.48, 1.28, 5.25, 2.95)
    add_card(
        slide,
        6.05,
        1.18,
        3.25,
        1.45,
        "Gradio 介面功能",
        "使用範例資料、上傳 CSV\n選擇分析欄位、測站、圖表\n顯示資料預覽、統計摘要與圖表",
        "cyan",
        14,
        10.5,
    )
    cmd_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.05), Inches(2.94), Inches(3.25), Inches(1.18))
    cmd_box.fill.solid()
    cmd_box.fill.fore_color.rgb = rgb("F3F6F8")
    cmd_box.line.color.rgb = rgb(COLORS["aqua"])
    add_textbox(
        slide,
        "docker build -t ocean-data-dashboard\n  -f docker/Dockerfile .\ndocker run -p 7860:7860\n  ocean-data-dashboard\nhttp://localhost:7860",
        6.18,
        3.04,
        3.0,
        0.9,
        8.2,
        "text",
        False,
        font_face=FONT_EN,
    )
    add_badge(slide, "Docker", 7.1, 4.35, 0.92, "blue")

    # 12
    slide = prs.slides.add_slide(blank)
    set_background(slide, COLORS["pale"])
    add_page_title(slide, "結論與未來改善方向", "完成從資料處理、分析視覺化到 Docker 部署的完整流程", 12)
    add_card(slide, 0.65, 1.18, 4.1, 2.35, "專題成果", "完成海洋環境資料讀取與清洗流程。\n完成統計分析與視覺化功能。\n建立 Gradio 互動式儀表板。\n完成 Docker 容器化部署。", "blue", 16, 12)
    add_card(slide, 5.25, 1.18, 4.1, 2.35, "未來改善", "串接真實海洋觀測 API。\n加入即時資料更新功能。\n增加浪高或風速異常警示。\n加入地圖視覺化顯示測站位置。", "green", 16, 12)
    closing = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(4.18), Inches(8.1), Inches(0.55))
    closing.fill.solid()
    closing.fill.fore_color.rgb = rgb(COLORS["navy"])
    closing.line.fill.background()
    add_textbox(slide, "本專題展現巨量資料與雲端運算課程的技術整合能力。", 1.12, 4.34, 7.75, 0.2, 15, "white", True, PP_ALIGN.CENTER)

    for idx, slide in enumerate(prs.slides, start=1):
        if idx != 1:
            add_page_badge(slide, idx)

    prs.save(OUTPUT_PATH)


def main():
    chart_paths = make_chart_assets()
    create_presentation(chart_paths)
    print(f"created={OUTPUT_PATH}")
    print(f"bytes={OUTPUT_PATH.stat().st_size}")


if __name__ == "__main__":
    main()
